import os
from docx import Document
import markdown
import fitz
from openpyxl import load_workbook
from pptx import Presentation

EXTENSION_TO_TYPE = {
    'pdf': 'pdf',
    'doc': 'docx', 'docx': 'docx',
    'md': 'md',
    'txt': 'txt',
    'xls': 'excel', 'xlsx': 'excel',
    'ppt': 'ppt', 'pptx': 'ppt'
}

class FileProcessor:
    def extract_content(self, file_path):
        """提取文件内容"""
        _, ext = os.path.splitext(file_path)
        ext = ext.lower().strip('.')
        
        file_type = EXTENSION_TO_TYPE.get(ext)
        if not file_type:
            raise ValueError(f"不支持的文件类型: {ext}")
            
        if file_type == 'pdf':
            return self.read_pdf(file_path)
        elif file_type == 'docx':
            return self.read_word(file_path)
        elif file_type == 'md':
            return self.read_markdown(file_path)
        elif file_type == 'excel':
            return self.read_excel(file_path)
        elif file_type == 'txt':
            return self.read_txt(file_path)
        elif file_type == 'ppt':
            return self.read_ppt(file_path)
        else:
            raise ValueError(f"未知的文件类型: {file_type}")

    def read_pdf(self, file_path):
        """读取 PDF 文件内容"""
        text = []
        with fitz.open(file_path) as doc:
            for page in doc:
                text.append(page.get_text())
        return "\n".join(text)

    def read_word(self, file_path):
        """读取 Word 文档内容"""
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def read_markdown(self, file_path):
        """读取 Markdown 文件内容"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    def read_excel(self, file_path):
        """读取 Excel 文件内容"""
        wb = load_workbook(file_path)
        text = []
        for sheet in wb:
            text.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                text.append("\t".join(str(cell) if cell is not None else "" for cell in row))
        return "\n".join(text)

    def read_txt(self, file_path):
        """读取文本文件内容"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    def read_ppt(self, file_path):
        """读取 PPT 文件内容"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text) 